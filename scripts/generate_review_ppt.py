from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Webnovel Architect: Project Second Review"
    subtitle.text = "A Neuro-Symbolic Approach to Dynamic Graph-RAG\nfor Serialized Fiction Audio Synthesis"

    # Slide 2: Agenda
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Agenda"
    tf = body_shape.text_frame
    tf.text = "1. System Architecture"
    p = tf.add_paragraph()
    p.text = "2. Explanation of Architecture Components"
    p = tf.add_paragraph()
    p.text = "3. Implementation Details"
    p = tf.add_paragraph()
    p.text = "4. Results and Discussion"

    # Slide 3: System Architecture Overview
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "System Architecture: The Neuro-Symbolic Switchboard"
    tf = body_shape.text_frame
    tf.text = "Resolving the Casting Paradox via DyG-RAG"
    
    p = tf.add_paragraph()
    p.text = "Event-Centric Architecture: Shifts from static document retrieval to a Dynamic Event Graph."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "The Switchboard Pattern: Decouples contextual semantic extraction from deterministic temporal reasoning."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Zero-Local-GPU Footprint: Routes heavy LLM tasks to remote endpoints, retaining instantaneous graph processing on the local CPU."
    p.level = 1

    # Slide 4: Explanation of Components (1/2)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "Architecture Components (1/2)"
    tf = body_shape.text_frame
    tf.text = "Layer 1: The Eye (Neural Extraction Interface)"
    p = tf.add_paragraph()
    p.text = "Leverages LiteLLM (Groq/Gemini) to perform complex entity recognition and capture emotional valences and relationship mapping."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Layer 2: The Brain (Symbolic Runtime)"
    p = tf.add_paragraph()
    p.text = "Maintains a NetworkX Directed Acyclic Graph (DAG) for narrative memory, resolving Temporal Hallucination present in standard RAG architectures."
    p.level = 1

    # Slide 5: Explanation of Components (2/2)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "Architecture Components (2/2)"
    tf = body_shape.text_frame
    
    tf.text = "Layer 3: The Director (Reasoning & Graduation)"
    p = tf.add_paragraph()
    p.text = "Utilizes PageRank centrality combined with Temporal Decay (λ = 0.15)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Employs the Debut Prominence Quotient (DPQ) to preemptively assign voices before characters reach full narrative prominence."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Layer 4: The Voice Broker (Audio Synthesis)"
    p = tf.add_paragraph()
    p.text = "Proactively assigns persistent TTS profiles (Kokoro ONNX for main cast, Edge-TTS fallback) closing the acoustic diarization gap."
    p.level = 1

    # Slide 6: Implementation Details
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "Implementation Details"
    tf = body_shape.text_frame
    
    tf.text = "Dual-Fallback Ingestion Pipeline: Robust combination of LiteLLM (primary) and spaCy (deterministic fallback) for unbroken execution."
    
    p = tf.add_paragraph()
    p.text = "Persistent In-Memory Graph: Utilizing JSON-backed NetworkX for lightweight logic mapping without database bloat."
    
    p = tf.add_paragraph()
    p.text = "Hysteresis Graduation State Machine: Strict thresholds prevent voice assignment thrashing (Upper bound = 0.15, Lower bound = 0.05)."

    p = tf.add_paragraph()
    p.text = "Extensive Audio Interfacing: Integrated chapter chunking, context caching for LLMs, and WebVTT syncing for visualization."

    # Slide 7: Results and Discussion - Extraction & Latency
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "Results: Precision and Execution Speed"
    tf = body_shape.text_frame
    
    tf.text = "Metric 1: Entity Extraction Efficacy"
    p = tf.add_paragraph()
    p.text = "LiteLLM Pipeline achieved 100% precision and recall (Character F1: 100%) on the benchmark."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "spaCy fallback severely restricted at 40% F1, proving the necessity of the neural layer for genre ambiguity."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Metric 2: Traversal Latency (Zero-GPU Simulation)"
    p = tf.add_paragraph()
    p.text = "Sub-4 millisecond performance (3.2 ms for PageRank + Decay logic) on fully loaded graphs of 1000 dense nodes utilizing purely standard CPU hardware."
    p.level = 1

    # Slide 8: Results and Discussion - Temporal Decay
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "Results: Eradicating Temporal Hallucination"
    tf = body_shape.text_frame
    
    tf.text = "Metric 3: Temporal Decay Ablation"
    p = tf.add_paragraph()
    p.text = "Evaluated the narrative dropoff for a transient character ('Mother') over five chapters."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Static Baseline (Vector RAG, λ=0.0): Retained erroneously high, persistent relevancy above threshold -> temporal hallucination."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "DyG-RAG (λ=0.15): Achieved mathematically monotonic decay, triggering voice release safely by Chapter 5."
    p.level = 1

    # Slide 9: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "Conclusion: The Casting Paradox Resolved"
    tf = body_shape.text_frame
    
    tf.text = "Synthesizing Success"
    p = tf.add_paragraph()
    p.text = "Successfully separated proactive semantic disambiguation (Neural) from temporal lifecycle tracking (Symbolic)."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Zero-Local GPU processing and the Switchboard adapter model prove that serialized audio-drama production is economically viable for independent creators."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Provides a highly adaptable engine addressing both the diarization gap and the 'LLM Tax' on consumer equipment."
    p.level = 1

    prs.save('Project_Second_Review.pptx')
    print("Presentation generated successfully at Project_Second_Review.pptx")

if __name__ == '__main__':
    create_presentation()
